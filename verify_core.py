"""驗證核心教學檔案處理套件。"""

from __future__ import annotations

import importlib
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


MODULES = {
    "docx": "python-docx",
    "openpyxl": "openpyxl",
    "pptx": "python-pptx",
    "pypdf": "pypdf",
    "fitz": "PyMuPDF",
    "reportlab": "reportlab",
    "PIL": "Pillow",
    "matplotlib": "matplotlib",
    "qrcode": "qrcode",
    "markitdown": "markitdown",
    "docx2pdf": "docx2pdf",
    "ocrmypdf": "ocrmypdf",
    "pypdfium2": "pypdfium2",
}


def find_word() -> Path | None:
    if sys.platform != "win32":
        return None

    candidates = [
        Path(r"C:\Program Files\Microsoft Office\Root\Office16\WINWORD.EXE"),
        Path(r"C:\Program Files (x86)\Microsoft Office\Root\Office16\WINWORD.EXE"),
    ]
    try:
        import winreg

        for hive in (winreg.HKEY_LOCAL_MACHINE, winreg.HKEY_CURRENT_USER):
            try:
                with winreg.OpenKey(
                    hive,
                    r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\WINWORD.EXE",
                ) as key:
                    candidates.insert(0, Path(winreg.QueryValue(key, None)))
            except OSError:
                pass
    except ImportError:
        pass

    return next((path for path in candidates if path.is_file()), None)


def find_tesseract() -> Path | None:
    command = shutil.which("tesseract")
    candidates = [
        Path(command) if command else None,
        Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
        Path(r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"),
    ]
    return next((path for path in candidates if path and path.is_file()), None)


def verify_external_tools() -> Path:
    word = find_word()
    if word is None:
        raise RuntimeError("Microsoft Word is required by the core docx2pdf workflow")

    tesseract = find_tesseract()
    if tesseract is None:
        raise RuntimeError("Tesseract OCR is required by the core OCRmyPDF workflow")

    result = subprocess.run(
        [str(tesseract), "--list-langs"],
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        env=os.environ.copy(),
    )
    languages = f"{result.stdout}\n{result.stderr}"
    missing = [language for language in ("eng", "osd", "chi_tra") if language not in languages]
    if result.returncode != 0 or missing:
        raise RuntimeError(f"Tesseract languages missing: {', '.join(missing) or 'unknown error'}")

    os.environ["PATH"] = f"{tesseract.parent}{os.pathsep}{os.environ.get('PATH', '')}"
    return tesseract


def run_smoke_tests() -> None:
    import fitz
    import matplotlib
    import ocrmypdf
    import qrcode
    from docx import Document
    from docx2pdf import convert as convert_docx_to_pdf
    from markitdown import MarkItDown
    from openpyxl import Workbook
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation
    from pptx.util import Inches
    from pypdf import PdfReader
    from reportlab.pdfgen import canvas

    matplotlib.use("Agg")
    from matplotlib import pyplot as plt

    with tempfile.TemporaryDirectory(prefix="teacher-file-core-") as temp_dir:
        root = Path(temp_dir)

        docx_path = root / "sample.docx"
        document = Document()
        document.add_paragraph("TEACHER FILE DOCX")
        document.save(docx_path)

        docx_pdf_path = root / "sample-docx.pdf"
        convert_docx_to_pdf(str(docx_path), str(docx_pdf_path))
        assert docx_pdf_path.is_file(), "docx2pdf did not create a PDF"
        assert len(PdfReader(docx_pdf_path).pages) == 1

        xlsx_path = root / "sample.xlsx"
        workbook = Workbook()
        workbook.active["A1"] = "TEACHER FILE XLSX"
        workbook.save(xlsx_path)

        pptx_path = root / "sample.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "TEACHER FILE PPTX"
        presentation.save(pptx_path)

        pdf_path = root / "sample.pdf"
        pdf_canvas = canvas.Canvas(str(pdf_path))
        pdf_canvas.drawString(72, 720, "TEACHER FILE PDF")
        pdf_canvas.save()
        assert len(PdfReader(pdf_path).pages) == 1
        with fitz.open(pdf_path) as pdf_document:
            assert pdf_document.page_count == 1
            pdf_document[0].get_pixmap(matrix=fitz.Matrix(0.5, 0.5))

        image_path = root / "sample.png"
        Image.new("RGB", (32, 32), "white").save(image_path)
        qrcode.make("https://example.com/teacher-file-toolkit").save(root / "qr.png")

        figure = plt.figure(figsize=(2, 1))
        plt.plot([0, 1], [0, 1])
        figure.savefig(root / "chart.png")
        plt.close(figure)

        converter = MarkItDown()
        expected_text = {
            docx_path: "TEACHER FILE DOCX",
            xlsx_path: "TEACHER FILE XLSX",
            pptx_path: "TEACHER FILE PPTX",
            pdf_path: "TEACHER FILE PDF",
        }
        for source_path, expected in expected_text.items():
            result = converter.convert(str(source_path))
            assert expected in result.text_content, f"MarkItDown failed for {source_path.suffix}"

        scan_pdf_path = root / "sample-scan.pdf"
        scan_image = Image.new("RGB", (1200, 500), "white")
        try:
            ocr_font = ImageFont.truetype(r"C:\Windows\Fonts\arial.ttf", 72)
        except OSError:
            ocr_font = ImageFont.load_default()
        ImageDraw.Draw(scan_image).text((80, 180), "OCR TEST 123", fill="black", font=ocr_font)
        scan_image.save(scan_pdf_path, "PDF", resolution=200.0)

        ocr_pdf_path = root / "sample-ocr.pdf"
        ocrmypdf.ocr(
            scan_pdf_path,
            ocr_pdf_path,
            language=["eng"],
            output_type="pdf",
            optimize=0,
            progress_bar=False,
        )
        with fitz.open(ocr_pdf_path) as ocr_document:
            ocr_text = " ".join(page.get_text() for page in ocr_document)
        assert "OCR" in ocr_text and "123" in ocr_text, "OCRmyPDF text smoke test failed"


def main() -> int:
    failures: list[tuple[str, str]] = []
    for module_name, package_name in MODULES.items():
        try:
            importlib.import_module(module_name)
        except Exception as exc:  # pragma: no cover - intended diagnostic path
            failures.append((package_name, f"{type(exc).__name__}: {exc}"))

    if failures:
        print(f"CORE_FAIL: {len(MODULES) - len(failures)}/{len(MODULES)}")
        for package_name, error in failures:
            print(f"- {package_name}: {error}")
        return 1

    try:
        verify_external_tools()
        run_smoke_tests()
    except Exception as exc:  # pragma: no cover - intended diagnostic path
        print(f"CORE_SMOKE_FAIL: {type(exc).__name__}: {exc}")
        return 1

    version = '.'.join(str(part) for part in sys.version_info[:3])
    print(f"CORE_OK: {len(MODULES)}/{len(MODULES)} imports and file smoke tests (Python {version})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
